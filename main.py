from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create the most comprehensive PPT yet with deep elaboration for each point
prs = Presentation()

def add_bullet_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # Make title larger and blue for beauty
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(13)
        p.font.name = "Calibri"
        if point.startswith("-"):
            p.level = 1
        if point.startswith("Example") or point.startswith("Case"):
            p.font.italic = True
        if point.startswith("Impact") or point.startswith("Role"):
            p.font.bold = True
    return slide

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Parameters and Variables of Management System in Industrial Management"
slide.placeholders[1].text = "A Comprehensive Exploration of Industrial Stability and Adaptability Factors"

def add_section_divider(title, subtitle=None):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    if subtitle:
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = subtitle
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    return slide

# Ultra-Detailed Slides with more info, visuals, and dividers
slides_content = [
    ("Section Divider", ["Section: Fundamentals of Industrial Management"]),
    ("Introduction", [
        "Industrial Management integrates engineering principles with management practices to optimize industrial operations.",
        "Key objectives include maximizing productivity, ensuring quality, minimizing costs, and maintaining safety.",
        "The 5 M's (Men, Materials, Machines, Methods, Money) are the core resources managed for efficiency.",
        "A Management System is a structured framework of policies, processes, and procedures for achieving organizational goals.",
        "Parameters and Variables are two fundamental concepts:",
        "- Parameters: Fixed, measurable limits or standards (e.g., machine capacity, safety norms).",
        "- Variables: Dynamic, changing factors (e.g., demand, workforce availability).",
        "Mastering both enables managers to balance stability with adaptability, ensuring long-term success."
        "[Visual Suggestion: Insert a flowchart showing the relationship between the 5 M's and management outcomes]"
    ]),
    ("Section Divider", ["Section: Deep Dive into Parameters"]),
    ("Understanding Parameters", [
        "Definition: Parameters are the set, measurable, and relatively constant constraints that define the boundaries of industrial processes.",
        "Role: Serve as the foundation for planning, standardization, and regulatory compliance.",
        "Characteristics:",
        "- Stable over time, rarely change unless there is a major policy or technological shift.",
        "- Quantifiable and often documented in SOPs or regulatory guidelines.",
        "- Aligned with industry standards and legal requirements.",
        "Impact:",
        "- Provide operational certainty and predictability.",
        "- Enable benchmarking and quality control.",
        "- Facilitate audits and certifications.",
        "Examples:",
        "- Maximum machine load (e.g., 10 tons per shift).",
        "- Standard cycle time (e.g., 2 minutes per unit).",
        "- ISO 9001 quality thresholds.",
        "- Regulatory environmental emission limits.",
        "- Minimum safety stock levels."
        "[Visual Suggestion: Table comparing different types of parameters in various industries]"
    ]),
    ("Section Divider", ["Section: Deep Dive into Variables"]),
    ("Understanding Variables", [
        "Definition: Variables are factors that fluctuate due to internal or external influences, impacting process outcomes and decision-making.",
        "Role: Introduce flexibility, allowing managers to adapt to real-time changes and uncertainties.",
        "Characteristics:",
        "- Time-sensitive and subject to frequent change.",
        "- Require ongoing monitoring and rapid response.",
        "- Often influenced by market, technology, or human factors.",
        "Impact:",
        "- Affect production efficiency, costs, and delivery timelines.",
        "- Can create opportunities or risks depending on how they are managed.",
        "Examples:",
        "- Sudden market demand surges or drops.",
        "- Employee absenteeism or turnover.",
        "- Supply chain disruptions (e.g., delayed shipments).",
        "- Raw material price volatility.",
        "- Adoption rate of new technologies."
        "[Visual Suggestion: Line graph showing variable trends over time]"
    ]),
    ("Section Divider", ["Section: Practical Applications"]),
    ("Key Parameters in Industrial Management", [
        "Quality Standards:",
        "- Benchmarks such as ISO, BIS, and Six Sigma ensure consistent output and customer satisfaction.",
        "- Example: ISO 9001 certification for quality management systems.",
        "Capacity Constraints:",
        "- Physical and operational limits of machinery, manpower, and facilities.",
        "- Example: Maximum daily output of a production line.",
        "Budget Limits:",
        "- Fixed cost allocations for procurement, production, and marketing.",
        "- Example: Annual maintenance budget for equipment.",
        "Safety Standards:",
        "- OSHA and other codes to ensure safe workplaces and prevent accidents.",
        "- Example: Mandated use of PPE in hazardous areas.",
        "Legal Compliance:",
        "- Adherence to labor laws, environmental acts, and trade regulations.",
        "- Example: Compliance with minimum wage laws.",
        "Operational Timelines:",
        "- Predefined schedules for projects or production runs.",
        "- Example: Delivery deadlines for customer orders."
        "[Visual Suggestion: Infographic of key parameters with icons]"
    ]),
    ("Key Variables in Industrial Management", [
        "Human Resource Availability:",
        "- Variations in skills, morale, turnover, and training levels.",
        "- Example: Sudden absenteeism during a pandemic.",
        "Market Trends:",
        "- Shifts in consumer preferences, seasonal demand, and competitor actions.",
        "- Example: Increased demand for eco-friendly products.",
        "Technological Changes:",
        "- Innovations like IoT, AI, and robotics transforming production.",
        "- Example: Automation reducing manual labor needs.",
        "Supply Chain Conditions:",
        "- Fluctuations in lead times, supplier reliability, and logistics.",
        "- Example: Delays due to geopolitical events.",
        "Energy Costs:",
        "- Changes in fuel and electricity prices affecting profitability.",
        "- Example: Rising energy costs during global crises.",
        "Economic and Political Climate:",
        "- Tariffs, inflation, and policy changes influencing operations.",
        "- Example: New import tariffs increasing raw material costs."
        "[Visual Suggestion: Pie chart of variable impact distribution]"
    ]),
    ("Relationship Between Parameters & Variables", [
        "Parameters provide stability and set the boundaries for operations; variables introduce adaptability and responsiveness.",
        "Their interaction determines how organizations balance efficiency with flexibility.",
        "Case Example:",
        "- A factory with a fixed 500-unit daily capacity (parameter) must adjust production schedules based on fluctuating demand (variable) to avoid overproduction or shortages.",
        "- During a supply chain disruption (variable), managers may need to operate within minimum inventory levels (parameter) to maintain production.",
        "Effective managers distinguish between non-negotiable limits and adjustable elements to optimize performance."
        "[Visual Suggestion: Venn diagram showing overlap of parameters and variables]"
    ]),
    ("Impact on Decision-Making", [
        "Parameters enable managers to set clear, achievable goals based on fixed resources and compliance needs.",
        "Variables require vigilance and agility to make timely adjustments in response to changing conditions.",
        "Balanced management prevents over-reliance on rigid rules or unpredictable factors.",
        "Ignoring parameters can cause compliance failures or inefficiencies; ignoring variables can lead to missed opportunities and competitive loss.",
        "Example: A company that ignores rising energy costs (variable) may exceed its budget limits (parameter), impacting profitability."
        "[Visual Suggestion: Decision tree for parameter-variable management]"
    ]),
    ("Case Study – Automobile Manufacturing", [
        "Parameters:",
        "- Fixed model blueprints and specifications.",
        "- Legally mandated safety features (e.g., airbags, crash tests).",
        "- Maximum assembly line speed (e.g., 60 cars/hour).",
        "- Pre-approved vendor lists for parts sourcing.",
        "Variables:",
        "- Fluctuating fuel and raw material prices.",
        "- Market demand shifts toward electric vehicles.",
        "- Labor disputes and workforce availability.",
        "- Supply shortages due to global events (e.g., chip shortage).",
        "Response Strategies:",
        "- Adjusting production volumes to match demand.",
        "- Diversifying suppliers to reduce risk.",
        "- Introducing flexible work shifts.",
        "- Realigning marketing campaigns to target new trends."
        "[Visual Suggestion: Timeline of case study events]"
    ]),
    ("Best Practices", [
        "Centralize documentation of all operational parameters for easy access and review.",
        "Use advanced analytics and ERP systems to monitor and predict variable changes in real time.",
        "Provide ongoing training for managers in adaptive decision-making and scenario planning.",
        "Conduct regular reviews of parameters to ensure they remain relevant and achievable.",
        "Develop contingency plans for high-impact variable changes, such as supply chain breakdowns or market crashes.",
        "Foster a culture of continuous improvement and feedback."
        "[Visual Suggestion: Checklist graphic for best practices]"
    ]),
    ("Conclusion", [
        "Industrial excellence is achieved by strategically blending fixed parameters with adaptive variables.",
        "Parameters ensure stability, compliance, and predictability; variables drive innovation, flexibility, and responsiveness.",
        "Organizations that master both can maintain operational balance and quickly adapt to new opportunities and threats.",
        "Continuous learning and improvement are key to sustaining success in a dynamic industrial environment."
        "[Visual Suggestion: Trophy or success icon]"
    ]),
    ("References", [
        "Harold Koontz & Heinz Weihrich – Essentials of Management",
        "Stevenson – Operations Management",
        "ISO 9001: Quality Management Standards",
        "Industrial Engineering & Management Journals",
        "OSHA Safety Guidelines",
        "Industry 4.0 Case Studies and White Papers",
        "Recent articles from Harvard Business Review and McKinsey on industrial adaptability"
    ])
]


# Add slides, using dividers for sections
for title, bullets in slides_content:
    if title == "Section Divider":
        add_section_divider(bullets[0])
    else:
        add_bullet_slide(title, bullets)

# Save file
pptx_path_ultra_detailed = "Parameters_Variables_Industrial_Management_Ultra_Detailed_03.pptx"
prs.save(pptx_path_ultra_detailed)

pptx_path_ultra_detailed
